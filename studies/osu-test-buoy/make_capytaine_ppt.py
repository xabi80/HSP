"""Build a small meeting deck explaining the Capytaine diffraction & radiation
analysis (and the pre-test heave-decay prediction) for the OSU Test Buoy.

Companion to OSU-TEST-BUOY-GEOMETRY.md; embeds Capytaine_explained.png and
OSU_heave_decay_prediction.png. Writes Capytaine_analysis_OSU_buoy.pptx next to
this script. Requires python-pptx (scripting-only, not a FloatSim runtime dep).
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
TEAL = RGBColor(0x0C, 0x8B, 0x96)
TEAL_D = RGBColor(0x0A, 0x55, 0x60)
INK = RGBColor(0x25, 0x32, 0x3A)
GREY = RGBColor(0x54, 0x63, 0x6D)
LIGHT = RGBColor(0xEE, 0xF6, 0xF7)
RED = RGBColor(0xD1, 0x54, 0x3A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Segoe UI"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def _bg(slide, color=WHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _rect(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def _title(slide, text, eyebrow=None):
    if eyebrow:
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.42), Inches(12), Inches(0.4))
        p = tb.text_frame.paragraphs[0]
        r = p.add_run(); r.text = eyebrow.upper()
        r.font.name = FONT; r.font.size = Pt(13); r.font.bold = True
        r.font.color.rgb = TEAL
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.72), Inches(12), Inches(0.9))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(30); r.font.bold = True
    r.font.color.rgb = INK
    _rect(slide, 0.72, 1.58, 2.2, 0.06, TEAL)


def _bullets(slide, items, l=0.85, t=1.95, w=11.7, h=5.0, size=18, gap=8):
    """items: list of (level, runs) where runs = list of (text, {opts})."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (lvl, runs) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        p.space_after = Pt(gap)
        p.space_before = Pt(2)
        bullet = "•  " if lvl == 0 else "–  "
        indent = "" if lvl == 0 else "      "
        lead = p.add_run(); lead.text = indent + bullet
        lead.font.name = FONT; lead.font.size = Pt(size)
        lead.font.color.rgb = TEAL if lvl == 0 else GREY
        for txt, opts in runs:
            r = p.add_run(); r.text = txt
            r.font.name = FONT
            r.font.size = Pt(opts.get("size", size if lvl == 0 else size - 2))
            r.font.bold = opts.get("b", False)
            r.font.italic = opts.get("i", False)
            r.font.color.rgb = opts.get("c", INK)
    return tb


def _panel(slide, l, t, w, h, header, sub, items, hc):
    """A titled column: colored header bar + subtitle + bullet list."""
    body = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    body.fill.solid(); body.fill.fore_color.rgb = RGBColor(0xF7, 0xFA, 0xFB)
    body.line.color.rgb = hc; body.line.width = Pt(1.25); body.shadow.inherit = False
    hb = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(0.6))
    hb.fill.solid(); hb.fill.fore_color.rgb = hc; hb.line.fill.background(); hb.shadow.inherit = False
    p = hb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = header
    r.font.name = FONT; r.font.size = Pt(17); r.font.bold = True; r.font.color.rgb = WHITE
    stb = slide.shapes.add_textbox(Inches(l), Inches(t + 0.58), Inches(w), Inches(0.34))
    p = stb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = sub
    r.font.name = FONT; r.font.size = Pt(10.5); r.font.italic = True; r.font.color.rgb = GREY
    tb = slide.shapes.add_textbox(Inches(l + 0.28), Inches(t + 1.02), Inches(w - 0.52), Inches(h - 1.15))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (txt, opts) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(opts.get("gap", 6))
        lead = p.add_run(); lead.text = opts.get("mark", "•  ")
        lead.font.name = FONT; lead.font.size = Pt(13); lead.font.color.rgb = hc
        r = p.add_run(); r.text = txt
        r.font.name = FONT; r.font.size = Pt(13)
        r.font.bold = opts.get("b", False)
        r.font.italic = opts.get("i", False)
        r.font.color.rgb = opts.get("c", INK)


def _cell(tbl, r, c, runs, fill, size=12, align=PP_ALIGN.LEFT):
    cell = tbl.cell(r, c)
    cell.fill.solid(); cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Pt(7); cell.margin_right = Pt(5)
    cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    if isinstance(runs, str):
        runs = [(runs, {})]
    for txt, opts in runs:
        rr = p.add_run(); rr.text = txt
        rr.font.name = FONT; rr.font.size = Pt(opts.get("size", size))
        rr.font.bold = opts.get("b", False); rr.font.italic = opts.get("i", False)
        rr.font.color.rgb = opts.get("c", INK)


def slide(title=None, eyebrow=None):
    s = prs.slides.add_slide(BLANK)
    _bg(s)
    if title:
        _title(s, title, eyebrow)
    return s


def T(text, **kw):
    return (text, kw)


# ---------------------------------------------------------------- 1. Title
s = slide()
_rect(s, 0, 0, 13.333, 7.5, WHITE)
_rect(s, 0, 5.55, 13.333, 1.95, TEAL)
tb = s.shapes.add_textbox(Inches(0.9), Inches(1.35), Inches(11.5), Inches(2.6))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run(); r.text = "DIFFRACTION & RADIATION ANALYSIS"
r.font.name = FONT; r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = TEAL
p2 = tf.add_paragraph(); p2.space_before = Pt(6)
r = p2.add_run(); r.text = "How we compute the buoy's\nhydrodynamics in Capytaine"
r.font.name = FONT; r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = INK
tb = s.shapes.add_textbox(Inches(0.95), Inches(5.75), Inches(11.5), Inches(1.5))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run(); r.text = "OSU Test Buoy  ·  potential-flow BEM  →  time-domain FloatSim (Cummins)"
r.font.name = FONT; r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = WHITE
p2 = tf.add_paragraph()
r = p2.add_run(); r.text = "Two sub-problems (radiation + diffraction), their outputs, and how they drive the simulation"
r.font.name = FONT; r.font.size = Pt(14); r.font.color.rgb = LIGHT

# ---------------------------------------------------------- 2. What it computes
s = slide("What Capytaine computes", "the tool")
_bullets(s, [
    (0, [T("Boundary-element (panel) solver for "), T("linear potential-flow hydrodynamics", b=True), T(".")]),
    (0, [T("In: the wetted hull as a mesh of panels.  Out: the flow (velocity potential φ) around it in regular waves at each frequency ω.")]),
    (0, [T("“Potential flow” = ", ), T("inviscid, irrotational, small (linear)", b=True, c=TEAL_D), T(" waves & motions — the one assumption that sets both the outputs and the limits.")]),
    (0, [T("Key idea — ", ), T("superposition", b=True), T(": the wave-body problem splits into two sub-problems, solved separately and added:")]),
    (1, [T("Radiation", b=True, c=TEAL_D), T(" — the body wavemaking in still water")]),
    (1, [T("Diffraction", b=True, c=TEAL_D), T(" — the body held fixed under incoming waves")]),
], t=2.0, gap=12)

# ---------------------------------------------------------- 3. BEM mesh
s = slide("The mesh we solve on", "the model")
s.shapes.add_picture(str(HERE / "OSU_buoy_mesh_capytaine.png"), Inches(3.35), Inches(1.5), height=Inches(4.98))
_rect(s, 0.85, 6.6, 11.6, 0.7, LIGHT)
tb = s.shapes.add_textbox(Inches(1.05), Inches(6.64), Inches(11.2), Inches(0.62))
tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
for txt, opts in [
    T("6″ spar cylinder + a "),
    T("placeholder solid equal-area disc", b=True, c=TEAL_D),
    T(" for the perforated heave plate.  Capytaine meshes the "),
    T("wetted part (z < 0, 3408 panels)", b=True),
    T("; the plate’s real hydro comes from the tank."),
]:
    r = p.add_run(); r.text = txt
    r.font.name = FONT; r.font.size = Pt(13); r.font.bold = opts.get("b", False)
    r.font.color.rgb = opts.get("c", INK)

# ---------------------------------------------------------- 4. Figure
s = slide("The two sub-problems — and how they combine", "overview")
pic = s.shapes.add_picture(str(HERE / "Capytaine_explained.png"), Inches(1.75), Inches(1.85), height=Inches(5.2))

# ---------------------------------------------------------- 4. Radiation
s = slide("Radiation — “body as a wavemaker”", "sub-problem ①")
_bullets(s, [
    (0, [T("Turn incident waves ", ), T("off", b=True), T(". Force the body to oscillate in still water — one DOF at a time, at each ω.")]),
    (0, [T("It radiates its own waves; the reaction force has two pieces:")]),
    (1, [T("Added mass A(ω)", b=True, c=TEAL_D), T(" — in phase with acceleration (virtual mass of entrained water).")]),
    (1, [T("Radiation damping B(ω)", b=True, c=TEAL_D), T(" — in phase with velocity (energy carried off as radiated waves).")]),
    (0, [T("6×6 matrices per ω — captures cross-DOF coupling (e.g. heave → pitch).")]),
    (0, [T("→ This run produces the ", ), T("A(ω) and B(ω)", b=True), T(" curves.")]),
], t=2.0, gap=12)

# ---------------------------------------------------------- 5. Diffraction
s = slide("Diffraction — “body as an obstacle”", "sub-problem ②")
_bullets(s, [
    (0, [T("Hold the body ", ), T("fixed", b=True), T(". Let incident waves scatter off it.")]),
    (0, [T("Net wave-pressure force on the fixed hull = ", ), T("excitation force F_exc(ω, β)", b=True, c=TEAL_D), T("  (β = wave heading):")]),
    (1, [T("Froude–Krylov", b=True), T(" — force from the undisturbed incident-wave pressure.")]),
    (1, [T("Diffraction", b=True), T(" — correction because the body scatters the waves.")]),
    (0, [T("→ This run produces the wave forcing ", ), T("F_exc(ω)", b=True), T(" that drives the motion.")]),
    (0, [T("Plus the hydrostatic stiffness ", ), T("C", b=True), T(" (buoyancy / waterplane) from the geometry.")]),
], t=2.0, gap=12)

# ---------------------------------------------------------- 6. Cummins bridge
s = slide("From frequency domain to time-domain simulation", "the bridge")
_bullets(s, [
    (0, [T("A(ω), B(ω) can’t go straight into a time-stepper — a transient contains all frequencies at once. Bridge = the "), T("Cummins equation", b=True), T(":")]),
], t=1.9, gap=6)
_rect(s, 0.85, 2.75, 11.6, 1.15, LIGHT)
tb = s.shapes.add_textbox(Inches(1.0), Inches(2.9), Inches(11.3), Inches(0.9))
tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "(M + A∞) ẍ  +  ∫₀ᵗ K(t−τ) ẋ(τ) dτ  +  C x  =  F_exc(t)  +  F_drag(ẋ)"
r.font.name = "Cambria Math"; r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = INK
_bullets(s, [
    (0, [T("A∞", b=True, c=TEAL_D), T(" = infinite-frequency added mass.")]),
    (0, [T("K(t) = (2/π)∫ B(ω) cos(ωt) dω", b=True, c=TEAL_D), T(" — the retardation kernel = ", ), T("fluid memory", i=True), T(" of previously radiated waves.")]),
    (0, [T("Mapping:  ", ), T("Radiation → A∞ + K(t)", b=True), T("   ·   "), T("Diffraction → F_exc", b=True), T("   ·   "), T("Hydrostatics → C", b=True), T(".  FloatSim integrates this in time.")]),
], t=4.1, gap=10)

# ------------------------------------------------ 8. Who computes what
s = slide("Who computes what — Capytaine vs FloatSim", "division of labour")
_panel(s, 0.55, 1.95, 5.5, 4.55, "CAPYTAINE",
       "run once · frequency domain · linear + inviscid", [
           T("Panel mesh of the wetted hull"),
           T("Radiation problem → A(ω), A∞ and B(ω)"),
           T("Diffraction problem → F_exc(ω, β)"),
           T("Hydrostatic stiffness C (buoyancy)"),
           T("Outputs: hydrodynamic coefficients", b=True, c=TEAL_D, mark="⇒  ", gap=2),
       ], TEAL)
_panel(s, 7.28, 1.95, 5.5, 4.55, "FLOATSIM",
       "per simulation · time domain", [
           T("Retardation kernel K(t) from B(ω)"),
           T("Body mass & inertia M (structure)"),
           T("Assemble (M+A∞), C + gravity m·g·z_G"),
           T("Viscous Morison drag F_drag — not in BEM", c=RED),
           T("Static equilibrium solve"),
           T("Time-integrate the Cummins EoM → x(t)"),
           T("Output: decay period T & damping ζ", b=True, c=TEAL_D, mark="⇒  ", gap=2),
       ], TEAL_D)
# hand-off arrow + label between the panels
arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.12), Inches(3.55), Inches(1.02), Inches(0.5))
arr.fill.solid(); arr.fill.fore_color.rgb = GREY; arr.line.fill.background(); arr.shadow.inherit = False
lab = s.shapes.add_textbox(Inches(5.75), Inches(4.08), Inches(1.75), Inches(0.7))
p = lab.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "capytaine_\nosu_buoy.nc"
r.font.name = "Consolas"; r.font.size = Pt(9); r.font.color.rgb = GREY
_rect(s, 0.85, 6.68, 11.95, 0.62, LIGHT)
tb = s.shapes.add_textbox(Inches(1.05), Inches(6.72), Inches(11.5), Inches(0.54))
tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
for txt, opts in [
    T("Capytaine = the hydrodynamic coefficients", b=True, c=TEAL_D),
    T("   ·   "),
    T("FloatSim = the equation of motion in time", b=True, c=TEAL_D),
    T("  (adds mass, viscous drag, mooring, equilibrium; converts B(ω)→K(t)).  "
      "Decay T and ζ are a FloatSim output."),
]:
    r = p.add_run(); r.text = txt
    r.font.name = FONT; r.font.size = Pt(12.5); r.font.bold = opts.get("b", False)
    r.font.color.rgb = opts.get("c", INK)

# ---------------------------------------------------------- 9. Limits
s = slide("Limits — and why the heave plate needs the tank", "caveats")
_bullets(s, [
    (0, [T("Potential flow is ", ), T("inviscid and linear", b=True), T(". Two consequences for us:")]),
    (1, [T("No viscous drag in the BEM", b=True, c=RED), T(" → all quadratic / viscous damping is added separately as ", ), T("Morison drag F_drag", b=True), T(" — not from Capytaine.")]),
    (1, [T("Panels treat every surface as solid", b=True, c=RED), T(" → ", ), T("over-predicts added mass", b=True), T(" for the open/perforated heave-plate frame (it blocks flow the real perforations pass).")]),
    (0, [T("Practical split for the OSU buoy:")]),
    (1, [T("Spar (6″ pipe)", b=True, c=TEAL_D), T(" — a clean cylinder → the BEM is reliable.")]),
    (1, [T("Heave-plate / ballast frame", b=True, c=TEAL_D), T(" — A and B are BEM ", ), T("stand-ins until the tank test measures them", b=True), T(".")]),
], t=2.0, gap=11)

# ------------------------------------------------ 10. Assumptions
s = slide("Modelling assumptions & their basis", "what we assumed, and why")
data = [
    ([T("Spar drag  "), T("Cd", i=True), T(" (transverse)")], [T("1.2", b=True)],
     [T("Smooth circular cylinder (Morison). Standard 1.0–1.2 (DNV-RP-C205; Sarpkaya). "),
      T("Sets surge & pitch drag; negligible in heave.", i=True)]),
    ([T("Heave-plate normal  "), T("Cd_n", i=True)], [T("5.0", b=True, c=RED)],
     [T("Flat disc at low KC (≈ 2–3 at a 100 mm release): flow separates, published Cd ≈ 4–8. "
        "Dominates heave damping.  "), T("Placeholder → the tank pins it.", b=True, c=RED)]),
    ([T("Heave-plate tangential  "), T("Cd_t", i=True)], [T("1.5", b=True)],
     [T("Edge / skin friction on the rim; small contribution.")]),
    ([T("Plate added mass")], [T("solid disc, Ø0.287 m", b=True, c=RED)],
     [T("Potential-flow UPPER bound — the perforated/webbed frame passes flow, so it adds "
        "less (porous-disk theory, Molin).  "), T("Bracket → the tank pins A33.", b=True, c=RED)]),
    ([T("Mass / CoG")], [T("21.52 kg;  −0.907 m", b=True)],
     [T("Spreadsheet structure 8.16 kg + measured unloaded waterline 967 mm. Independent "
        "geometry check (cyl + frame + lead) agrees within ~1.5%.")]),
    ([T("Pitch / roll inertia")], [T("10.2 kg·m²", b=True)],
     [T("gmsh per-part inertia tensors + uniform effective density; the deep lead ballast "
        "dominates, so it is robust to the internal mass split. CATIA would refine.")]),
    ([T("Water & BEM")], [T("fresh; deep; linear", b=True)],
     [T("OSU Hinsdale lab is fresh (ρ = 998); period is density-independent for a free body. "
        "Potential flow is inviscid — viscous drag added via Morison. Confirm tank depth vs 1.42 m draft.")]),
]
tbl = s.shapes.add_table(len(data) + 1, 3, Inches(0.55), Inches(1.82),
                         Inches(12.25), Inches(4.9)).table
tbl.first_row = False; tbl.horz_banding = False
tbl.columns[0].width = Inches(3.05)
tbl.columns[1].width = Inches(2.35)
tbl.columns[2].width = Inches(6.85)
for c, head in enumerate(["Assumption", "Value", "Basis / support"]):
    _cell(tbl, 0, c, [T(head, b=True, c=WHITE)], TEAL, size=13,
          align=PP_ALIGN.CENTER if c == 1 else PP_ALIGN.LEFT)
tbl.rows[0].height = Inches(0.42)
for i, (a, v, bss) in enumerate(data, start=1):
    fill = RGBColor(0xF3, 0xF8, 0xF9) if i % 2 else WHITE
    _cell(tbl, i, 0, a, fill, size=12)
    _cell(tbl, i, 1, v, fill, size=12, align=PP_ALIGN.CENTER)
    _cell(tbl, i, 2, bss, fill, size=11)
    tbl.rows[i].height = Inches(0.64)
tb = s.shapes.add_textbox(Inches(0.55), Inches(6.86), Inches(12.25), Inches(0.5))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
for txt, opts in [
    T("Red = placeholder the tank measures", b=True, c=RED),
    T("  (period → plate added mass, decay rate → plate Cd).  Coefficient ranges per "
      "DNV-RP-C205 & Sarpkaya; perforated-plate reduction per porous-disk theory.", c=GREY),
]:
    r = p.add_run(); r.text = txt
    r.font.name = FONT; r.font.size = Pt(10.5); r.font.bold = opts.get("b", False)
    r.font.color.rgb = opts.get("c", GREY)

# --------------------------------------------------------- 11. Prediction
s = slide("Pre-test heave decay prediction", "what to expect")
s.shapes.add_picture(str(HERE / "OSU_heave_decay_prediction.png"), Inches(1.15), Inches(1.85), width=Inches(11.0))
_rect(s, 0.85, 6.35, 11.6, 0.82, LIGHT)
tb = s.shapes.add_textbox(Inches(1.05), Inches(6.4), Inches(11.2), Inches(0.72))
tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
for txt, opts in [
    T("Period is predictable → T ≈ 2.3–2.4 s", b=True, c=TEAL_D),
    T("   (M, C₃₃ known; only plate added mass uncertain).   ", ),
    T("Damping is the measurement", b=True, c=RED),
    T("  → ζ₁ ≈ 8–15% at 100 mm.  Test pins: period→added mass, decay→drag Cd.", ),
]:
    r = p.add_run(); r.text = txt
    r.font.name = FONT; r.font.size = Pt(14); r.font.bold = opts.get("b", False)
    r.font.color.rgb = opts.get("c", INK)

out = HERE / "Capytaine_analysis_OSU_buoy.pptx"
prs.save(str(out))
print(f"wrote {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
